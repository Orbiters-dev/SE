"""
PPT 기획안: PPSU 200ml vs ワンタッチ 300ml 비교 카루셀 (5장)
인스타그램 카루셀 포맷 (1080x1080 px, 1:1 정사각형)
Reference: https://www.instagram.com/p/DXOCQU0ERo2/ (Grosmimi SG)
"""
import io
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

# 1080 px = 11.25 inches @ 96 DPI = 10,287,000 EMU
SIDE = Emu(10287000)

PINK = RGBColor(0xFF, 0x6B, 0x9D)
LIGHT_PINK = RGBColor(0xFF, 0xD6, 0xE0)
IVORY = RGBColor(0xFF, 0xF5, 0xF7)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xFF, 0x8A, 0x7A)

FONT_JP = "Yu Gothic UI"
FONT_BOLD = "Yu Gothic UI"


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SIDE, SIDE)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *, size=40, bold=False,
             color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_JP):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_label(slide, left, top, width, height, text, *, bg=PINK, fg=WHITE, size=24):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(100000)
    tf.margin_top = tf.margin_bottom = Emu(50000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape


def add_placeholder_photo(slide, left, top, width, height, label):
    """실제 제품 사진 들어갈 자리 (회색 박스 + 안내 텍스트)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_PINK
    shape.line.color.rgb = PINK
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[{label}]"
    run.font.name = FONT_JP
    run.font.size = Pt(22)
    run.font.color.rgb = PINK
    run.font.bold = True
    return shape


def add_two_color_title(slide, left, top, width, height, parts):
    """[('normal', '普通テキスト'), ('accent', 'PPSU'), ...] 식으로 키워드만 색칠"""
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for kind, text in parts:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_JP
        run.font.bold = True
        if kind == "accent":
            run.font.size = Pt(72)
            run.font.color.rgb = PINK
        elif kind == "accent_sm":
            run.font.size = Pt(56)
            run.font.color.rgb = PINK
        elif kind == "big":
            run.font.size = Pt(56)
            run.font.color.rgb = DARK
        else:
            run.font.size = Pt(44)
            run.font.color.rgb = DARK
    return tb


def slide1_cover(prs):
    """커버: 소제목(라벨) + 대제목(키워드 색 강조) + 두 제품 자리"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    # 소제목 라벨 (상단)
    add_label(s, 3543000, 800000, 3200000, 600000, "成長に合わせて選ぶ", bg=PINK, fg=WHITE, size=22)

    # 대제목 2줄
    add_two_color_title(s, 600000, 1600000, 9087000, 900000, [
        ("accent", "PPSU "), ("big", "200ml と"),
    ])
    add_two_color_title(s, 600000, 2500000, 9087000, 900000, [
        ("accent", "ワンタッチ "), ("big", "300ml"),
    ])
    add_text(s, 600000, 3500000, 9087000, 700000, "どう使い分ける？",
             size=48, bold=True, color=CORAL)

    # 두 제품 사진 placeholder (좌/우)
    add_placeholder_photo(s, 700000, 4600000, 4000000, 4000000, "PPSU 200ml\n실사진")
    add_placeholder_photo(s, 5500000, 4600000, 4000000, 4000000, "ワンタッチ 300ml\n실사진")

    # 하단 서명
    add_text(s, 0, 9600000, SIDE, 400000, "@grosmimi_japan",
             size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide2_ppsu(prs):
    """PPSU 200ml 상세"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    # 헤더
    add_label(s, 700000, 600000, 2400000, 500000, "1", bg=PINK, fg=WHITE, size=28)
    add_text(s, 3200000, 600000, 6000000, 500000, "PPSU 200ml",
             size=44, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    # 사진 placeholder
    add_placeholder_photo(s, 700000, 1400000, 4200000, 4200000, "PPSU 200ml\n제품 실사진")

    # 오른쪽 스펙 박스
    add_text(s, 5200000, 1500000, 4400000, 500000, "月齢の目安",
             size=22, bold=True, color=PINK, align=PP_ALIGN.LEFT)
    add_text(s, 5200000, 2000000, 4400000, 600000, "6ヶ月〜",
             size=42, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_text(s, 5200000, 2900000, 4400000, 500000, "素材",
             size=22, bold=True, color=PINK, align=PP_ALIGN.LEFT)
    add_text(s, 5200000, 3400000, 4400000, 500000, "医療用グレード PPSU",
             size=26, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_text(s, 5200000, 4100000, 4400000, 500000, "シーン",
             size=22, bold=True, color=PINK, align=PP_ALIGN.LEFT)
    add_text(s, 5200000, 4600000, 4400000, 500000, "おうち使い・入門に",
             size=26, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    # 하단 포인트 (체크마크)
    points = [
        "✔ 軽くて赤ちゃんも持ちやすい",
        "✔ パーツ交換でストローマグに変身",
        "✔ 煮沸・電子レンジ消毒OK",
    ]
    for i, pt in enumerate(points):
        add_text(s, 700000, 6300000 + i * 700000, 9000000, 600000, pt,
                 size=26, bold=False, color=DARK, align=PP_ALIGN.LEFT)

    add_text(s, 0, 9600000, SIDE, 400000, "2/5", size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide3_onetouch(prs):
    """ワンタッチ 300ml 상세"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    add_label(s, 700000, 600000, 2400000, 500000, "2", bg=CORAL, fg=WHITE, size=28)
    add_text(s, 3200000, 600000, 6000000, 500000, "ワンタッチ 300ml",
             size=44, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_placeholder_photo(s, 700000, 1400000, 4200000, 4200000, "ワンタッチ 300ml\n제품 실사진")

    add_text(s, 5200000, 1500000, 4400000, 500000, "月齢の目安",
             size=22, bold=True, color=CORAL, align=PP_ALIGN.LEFT)
    add_text(s, 5200000, 2000000, 4400000, 600000, "12ヶ月〜",
             size=42, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_text(s, 5200000, 2900000, 4400000, 500000, "特徴",
             size=22, bold=True, color=CORAL, align=PP_ALIGN.LEFT)
    add_text(s, 5200000, 3400000, 4400000, 500000, "ワンタッチ開閉・漏れない",
             size=26, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_text(s, 5200000, 4100000, 4400000, 500000, "シーン",
             size=22, bold=True, color=CORAL, align=PP_ALIGN.LEFT)
    add_text(s, 5200000, 4600000, 4400000, 500000, "お出かけ・保育園に",
             size=26, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    points = [
        "✔ 片手でサッと開けられる",
        "✔ カバン入れても漏れ知らず",
        "✔ 300mlで水分補給もたっぷり",
    ]
    for i, pt in enumerate(points):
        add_text(s, 700000, 6300000 + i * 700000, 9000000, 600000, pt,
                 size=26, bold=False, color=DARK, align=PP_ALIGN.LEFT)

    add_text(s, 0, 9600000, SIDE, 400000, "3/5", size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide4_compare(prs):
    """비교표"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    add_label(s, 3543000, 600000, 3200000, 500000, "くらべてみた", bg=PINK, fg=WHITE, size=22)
    add_text(s, 600000, 1300000, 9087000, 700000, "ここが違う！",
             size=48, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # 비교표
    rows = [
        ("項目", "PPSU 200ml", "ワンタッチ 300ml"),
        ("月齢", "6ヶ月〜", "12ヶ月〜"),
        ("容量", "200ml", "300ml"),
        ("フタ", "ストロー/スパウト", "ワンタッチ開閉"),
        ("素材", "PPSU(医療用)", "PPSU+ステンレス"),
        ("シーン", "おうち", "お出かけ"),
    ]
    table_top = 2400000
    row_h = 900000
    col_x = [500000, 3200000, 6300000]
    col_w = [2700000, 3100000, 3100000]

    for i, row in enumerate(rows):
        is_header = (i == 0)
        for j, cell in enumerate(row):
            shape = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(col_x[j]), Emu(table_top + i * row_h),
                Emu(col_w[j]), Emu(row_h),
            )
            shape.fill.solid()
            if is_header:
                shape.fill.fore_color.rgb = PINK
                text_color = WHITE
            elif i % 2 == 1:
                shape.fill.fore_color.rgb = WHITE
                text_color = DARK
            else:
                shape.fill.fore_color.rgb = LIGHT_PINK
                text_color = DARK
            shape.line.color.rgb = PINK
            shape.line.width = Pt(1)
            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Emu(100000)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = cell
            run.font.name = FONT_JP
            run.font.size = Pt(20 if is_header else 18)
            run.font.bold = is_header or (j == 0)
            run.font.color.rgb = text_color

    add_text(s, 0, 9600000, SIDE, 400000, "4/5", size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide5_conclusion(prs):
    """결론 + CTA"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    add_label(s, 3543000, 600000, 3200000, 500000, "おすすめの使い方", bg=PINK, fg=WHITE, size=22)

    add_two_color_title(s, 500000, 1500000, 9287000, 1200000, [
        ("big", "実は、"), ("accent", "両方使う"), ("big", "のが正解！"),
    ])

    # 두 제품 병행 이미지 placeholder
    add_placeholder_photo(s, 1500000, 3100000, 3400000, 3400000, "PPSU 200ml\n(おうち)")
    add_placeholder_photo(s, 5400000, 3100000, 3400000, 3400000, "ワンタッチ 300ml\n(お出かけ)")

    add_text(s, 500000, 6800000, 9287000, 600000, "おうちでPPSU、お出かけはワンタッチ。",
             size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 500000, 7500000, 9287000, 600000, "シーンで使い分けるのが賢い選択✨",
             size=24, bold=False, color=DARK, align=PP_ALIGN.CENTER)

    # CTA 박스
    cta = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(2000000), Emu(8500000), Emu(6287000), Emu(700000))
    cta.fill.solid()
    cta.fill.fore_color.rgb = PINK
    cta.line.fill.background()
    tf = cta.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "プロフィールからチェック 👆"
    run.font.name = FONT_JP
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_text(s, 0, 9600000, SIDE, 400000, "5/5 · @grosmimi_japan",
             size=16, color=GRAY, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SIDE
    prs.slide_height = SIDE

    slide1_cover(prs)
    slide2_ppsu(prs)
    slide3_onetouch(prs)
    slide4_compare(prs)
    slide5_conclusion(prs)

    out_dir = Path(r"C:\Users\orbit\Desktop\s\요청하신 자료\인스타그램 기획안")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "PPSU200_vs_ワンタッチ300_기획안.pptx"
    prs.save(out)
    print(f"[OK] PPT saved: {out}")


if __name__ == "__main__":
    main()
