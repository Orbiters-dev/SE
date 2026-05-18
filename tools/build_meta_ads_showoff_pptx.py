"""
ORBI Grosmimi JP — Meta Ads 운영 체계 (Show off PPT, 2026-05-18)

외부 노출용. 깔끔한 디자인 + 한 슬라이드 = 한 메시지.
세은이 외부 발표·공유 시 사용.

저장: C:\\Users\\orbit\\Desktop\\s\\메타\\08_meta_ads_showoff.pptx
"""
import sys
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT = Path(r"C:\Users\orbit\Desktop\s\메타\08_meta_ads_showoff.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

# 16:9 (default)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# Color palette
NAVY = RGBColor(0x1F, 0x29, 0x37)        # 본문 텍스트
ACCENT = RGBColor(0x25, 0x63, 0xEB)      # 파랑 액센트
ACCENT2 = RGBColor(0xDB, 0x27, 0x77)     # 핑크 액센트 (Grosmimi 톤)
LIGHT = RGBColor(0xF9, 0xFA, 0xFB)       # 배경
GRAY = RGBColor(0x6B, 0x72, 0x80)        # 부연
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_text(slide, x, y, w, h, text, font_size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Malgun Gothic"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    # 첫 줄
    if "\n" in text:
        lines = text.split("\n")
    else:
        lines = [text]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_rect(slide, x, y, w, h, fill=ACCENT, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is None:
        shape.line.fill.background()
    return shape


def add_header_bar(slide, page_num, total, section, title):
    # 좌상단 section + page indicator
    add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.4),
             section, font_size=12, color=GRAY, bold=True)
    add_text(slide, Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.4),
             f"{page_num} / {total}", font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    # 타이틀
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.9),
             title, font_size=32, bold=True, color=NAVY)
    # 액센트 라인
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(0.8), Emu(40000), fill=ACCENT)


TOTAL = 13


# ============================================================
# Slide 1 — 표지
# ============================================================
s = add_blank_slide()
# 배경
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
# 액센트 사각형 (장식)
add_rect(s, Inches(0), Inches(6.0), SLIDE_W, Inches(0.15), fill=ACCENT)
add_rect(s, Inches(9.5), Inches(0), Inches(0.15), SLIDE_H, fill=ACCENT2)

add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(0.6),
         "ORBI / Grosmimi JP", font_size=22, color=WHITE)
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(1.2),
         "Meta Ads", font_size=60, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(1.0),
         "운영 체계 2026", font_size=44, bold=True, color=ACCENT)

add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "2026-05-18 · act_4117678028561958",
         font_size=12, color=GRAY)


# ============================================================
# Slide 2 — 한 줄 요약
# ============================================================
s = add_blank_slide()
add_header_bar(s, 2, TOTAL, "WHY", "한 줄로 요약하면")

add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
         "KPI는 북극성, Off Rule은 동적 임계.",
         font_size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8),
         "두 룰을 섞으면 풀 평균이 KPI 미달이라 거의 모든 광고가 Off 대상이 된다.",
         font_size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
         "→ 분리해서, 풀 자체 percentile로만 판정한다.",
         font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 3 — 4개 원칙
# ============================================================
s = add_blank_slide()
add_header_bar(s, 3, TOTAL, "PRINCIPLES", "운영 4대 원칙")

principles = [
    ("01", "KPI ≠ Off Rule", "KPI는 고정 목표,\nOff Rule은 풀 percentile 기반 동적 임계.", ACCENT),
    ("02", "외부 시장 평균 폐기", "Triple Whale 1.91% 등\n시장 평균 직접 사용 X.", ACCENT2),
    ("03", "풀 분리 + 게이트", "PAID/GIFT × Testing/Proven\n4-way + newborn 동결.", GREEN),
    ("04", "Creative ≠ Budget", "두 영역은 독립.\n동시 진행 가능.", NAVY),
]
for i, (num, title, body, color) in enumerate(principles):
    col = i % 2
    row = i // 2
    x = Inches(0.7 + col * 6.3)
    y = Inches(2.2 + row * 2.4)
    add_rect(s, x, y, Inches(5.9), Inches(2.1), fill=LIGHT)
    add_rect(s, x, y, Inches(0.15), Inches(2.1), fill=color)
    add_text(s, x + Inches(0.4), y + Inches(0.15), Inches(1.2), Inches(0.6),
             num, font_size=32, bold=True, color=color)
    add_text(s, x + Inches(1.8), y + Inches(0.25), Inches(4), Inches(0.5),
             title, font_size=20, bold=True, color=NAVY)
    add_text(s, x + Inches(1.8), y + Inches(0.95), Inches(4), Inches(1.1),
             body, font_size=14, color=GRAY)


# ============================================================
# Slide 4 — KPI 북극성
# ============================================================
s = add_blank_slide()
add_header_bar(s, 4, TOTAL, "KPI", "북극성 (고정 목표)")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6),
         "달성 / 미달 카운터로만 운용. 평균 비교 X.",
         font_size=18, color=GRAY)

# 3개 KPI 박스
kpis = [
    ("WL CTR", "≥ 8.0%", "인플루언서\n(PAID + GIFT)", ACCENT),
    ("이미지 CTR", "≥ 4.0%", "자체 이미지\n광고", ACCENT2),
    ("CPC", "≤ ₩100", "전 광고\n공통", GREEN),
]
for i, (name, val, sub, color) in enumerate(kpis):
    x = Inches(0.7 + i * 4.3)
    y = Inches(3.0)
    add_rect(s, x, y, Inches(4.0), Inches(3.5), fill=LIGHT)
    add_rect(s, x, y, Inches(4.0), Inches(0.6), fill=color)
    add_text(s, x, y + Inches(0.05), Inches(4.0), Inches(0.5),
             name, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.0), Inches(4.0), Inches(1.5),
             val, font_size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.6), Inches(4.0), Inches(0.8),
             sub, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 5 — Off Rule 4축
# ============================================================
s = add_blank_slide()
add_header_bar(s, 5, TOTAL, "OFF RULE", "동적 임계 — 4축")

axes = [
    ("축 1", "풀 P20 미달", "풀 percentile P20 이하\n+ spend ≥ ₩30K", "메인 트리거", ACCENT),
    ("축 2", "시간 sunset", "D+60 이상\n(자연 노화)", "Off + 신규 변형 의뢰", ACCENT2),
    ("축 3", "피로도", "Freq ≥ 3.0\n(spend 무관)", "즉시 Off", RED),
    ("축 4", "velocity 하락", "Proven D+15+\ntrailing 7d CTR Δ < -15%", "신규 변형 우선", GRAY),
]
for i, (label, name, cond, action, color) in enumerate(axes):
    col = i % 2
    row = i // 2
    x = Inches(0.7 + col * 6.3)
    y = Inches(2.0 + row * 2.6)
    add_rect(s, x, y, Inches(5.9), Inches(2.3), fill=LIGHT)
    add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(1.5), Inches(0.5),
             label, font_size=14, bold=True, color=color)
    add_text(s, x + Inches(0.3), y + Inches(0.6), Inches(5.3), Inches(0.6),
             name, font_size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(1.25), Inches(5.3), Inches(0.7),
             cond, font_size=13, color=GRAY)
    add_text(s, x + Inches(0.3), y + Inches(1.85), Inches(5.3), Inches(0.4),
             f"→ {action}", font_size=13, bold=True, color=color)


# ============================================================
# Slide 6 — 풀 분리 (PAID/GIFT × Testing/Proven)
# ============================================================
s = add_blank_slide()
add_header_bar(s, 6, TOTAL, "POOL SPLIT", "4-way 풀 분리")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6),
         "회수 압력 · 학습 단계가 다른 광고를 한 풀로 평가하면 진단 정확도 하락.",
         font_size=16, color=GRAY)

# 2x2 매트릭스
mat_x = Inches(2.5)
mat_y = Inches(2.8)
cell_w = Inches(4.0)
cell_h = Inches(1.8)

# 헤더 (Testing / Proven)
add_text(s, mat_x + cell_w, mat_y - Inches(0.5), cell_w, Inches(0.4),
         "Testing (D+7~D+14)", font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, mat_x + cell_w*2, mat_y - Inches(0.5), cell_w, Inches(0.4),
         "Proven (D+15+)", font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
# 좌측 (PAID / GIFT)
add_text(s, mat_x - Inches(1.0), mat_y + Inches(0.6), Inches(1.0), Inches(0.4),
         "PAID", font_size=14, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
add_text(s, mat_x - Inches(1.0), mat_y + cell_h + Inches(0.6), Inches(1.0), Inches(0.4),
         "GIFT", font_size=14, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

# 셀
cells = [
    (0, 0, "PAID × Testing", "졸업 후보\n발굴", ACCENT2),
    (0, 1, "PAID × Proven", "Off 후보\n즉시", RED),
    (1, 0, "GIFT × Testing", "데이터\n누적", GRAY),
    (1, 1, "GIFT × Proven", "PAID 캐스팅\n전환", GREEN),
]
for row, col, title, body, color in cells:
    x = mat_x + cell_w * (col + 1)
    y = mat_y + cell_h * row
    add_rect(s, x, y, cell_w, cell_h, fill=LIGHT)
    add_rect(s, x, y, Inches(0.1), cell_h, fill=color)
    add_text(s, x + Inches(0.25), y + Inches(0.2), cell_w - Inches(0.4), Inches(0.5),
             title, font_size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.85), cell_w - Inches(0.4), Inches(0.8),
             body, font_size=12, color=GRAY)


# ============================================================
# Slide 7 — Lifecycle Gate
# ============================================================
s = add_blank_slide()
add_header_bar(s, 7, TOTAL, "LIFECYCLE", "변경 동결 게이트")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6),
         "광고 등록 후 시간에 따라 허용 액션이 다름. 학습 리셋 방지.",
         font_size=16, color=GRAY)

stages = [
    ("newborn", "D+0~D+6", "조회만\n변경 X", RED),
    ("testing", "D+7~D+14", "spend ≥ ₩30K 후\n평가", ACCENT2),
    ("proven", "D+15~D+59", "정상 운영\nKPI + Off + S 평가", GREEN),
    ("sunset", "D+60+", "Off 자동\n트리거 (축 2)", GRAY),
]
y0 = Inches(3.2)
for i, (name, period, action, color) in enumerate(stages):
    x = Inches(0.7 + i * 3.15)
    add_rect(s, x, y0, Inches(3.0), Inches(3.0), fill=LIGHT)
    add_rect(s, x, y0, Inches(3.0), Inches(0.7), fill=color)
    add_text(s, x, y0 + Inches(0.1), Inches(3.0), Inches(0.5),
             name, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.0), Inches(3.0), Inches(0.6),
             period, font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.9), Inches(3.0), Inches(1.0),
             action, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 8 — S 후보 발굴 (Top 10%)
# ============================================================
s = add_blank_slide()
add_header_bar(s, 8, TOTAL, "S 후보", "Top 10% 발굴")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6),
         "Motion 5% winner 룰을 풀 size 보정해 10%로 보수 설정.",
         font_size=16, color=GRAY)

# 좌측: 조건
add_text(s, Inches(0.7), Inches(2.8), Inches(5.5), Inches(0.5),
         "조건 (모두 충족)", font_size=18, bold=True, color=NAVY)
conds = [
    "• 풀 P90 이상 (Top 10%)",
    "• spend ≥ ₩30K (학습 통과)",
    "• 풀 n ≥ 6 (percentile 의미)",
    "• lifecycle = Proven",
    "• Off 트리거 hit 제외 (dedup)",
]
for i, c in enumerate(conds):
    add_text(s, Inches(0.7), Inches(3.5 + i*0.5), Inches(5.5), Inches(0.5),
             c, font_size=15, color=NAVY)

# 우측: 도식 (피라미드)
pyr_x = Inches(7.5)
pyr_y = Inches(3.0)
add_rect(s, pyr_x + Inches(2.0), pyr_y, Inches(1.5), Inches(0.6), fill=ACCENT2)
add_text(s, pyr_x + Inches(2.0), pyr_y + Inches(0.1), Inches(1.5), Inches(0.4),
         "S 후보 Top 10%", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, pyr_x + Inches(1.0), pyr_y + Inches(0.7), Inches(3.5), Inches(1.0), fill=GREEN)
add_text(s, pyr_x + Inches(1.0), pyr_y + Inches(1.0), Inches(3.5), Inches(0.4),
         "Middle 70%", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, pyr_x, pyr_y + Inches(1.8), Inches(5.5), Inches(1.0), fill=RED)
add_text(s, pyr_x, pyr_y + Inches(2.1), Inches(5.5), Inches(0.4),
         "Bottom 20% Off 후보", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 9 — Creative 표준 (5부분 구조)
# ============================================================
s = add_blank_slide()
add_header_bar(s, 9, TOTAL, "CREATIVE", "영상 광고 5부분 구조")

parts = [
    ("Hook", "첫 3초", "Shocking · Wow moment", ACCENT2, True),
    ("Education", "Hook 직후", "의학·발달 지식", ACCENT, True),
    ("Benefit", "중반", "제품 쓰면 좋은 점", GREEN, True),
    ("Result", "후반", "실제 사용 후 변화", GRAY, False),
    ("Ending", "마지막", "마무리 (CTA)", GRAY, False),
]
y0 = Inches(2.2)
for i, (name, pos, body, color, required) in enumerate(parts):
    y = y0 + Inches(i * 0.95)
    add_rect(s, Inches(0.7), y, Inches(11.9), Inches(0.85), fill=LIGHT)
    add_rect(s, Inches(0.7), y, Inches(0.15), Inches(0.85), fill=color)
    add_text(s, Inches(1.1), y + Inches(0.18), Inches(2.5), Inches(0.5),
             name, font_size=20, bold=True, color=NAVY)
    add_text(s, Inches(3.7), y + Inches(0.25), Inches(2.5), Inches(0.4),
             pos, font_size=13, color=GRAY)
    add_text(s, Inches(6.2), y + Inches(0.25), Inches(5.0), Inches(0.4),
             body, font_size=14, color=NAVY)
    if required:
        add_text(s, Inches(11.2), y + Inches(0.25), Inches(1.3), Inches(0.4),
                 "필수", font_size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 10 — 바이럴 의뢰 전략
# ============================================================
s = add_blank_slide()
add_header_bar(s, 10, TOTAL, "VIRAL STRATEGY", "바이럴 영상 의뢰 전략")

steps = [
    ("1", "리서치", "TikTok / Instagram\n관련 분야 바이럴 영상"),
    ("2", "평가", "좋아요·뷰 X\n저장 + 공유 + 코멘트"),
    ("3", "선별", "제품 자체 질문이\n코멘트에 많은 영상"),
    ("4", "동시 의뢰", "참고 hook을 인플루언서\n10명한테 동시 발주"),
    ("5", "바이럴", "한 명 뜨면\n같은 풀 전체 효과"),
]
y0 = Inches(2.5)
for i, (num, title, body) in enumerate(steps):
    x = Inches(0.5 + i * 2.55)
    # 동그라미 번호
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y0, Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    add_text(s, x + Inches(0.8), y0 + Inches(0.1), Inches(0.7), Inches(0.5),
             num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.0), Inches(2.3), Inches(0.5),
             title, font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.6), Inches(2.3), Inches(1.5),
             body, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
    # 화살표 (마지막 제외)
    if i < 4:
        add_text(s, x + Inches(2.2), y0 + Inches(0.15), Inches(0.5), Inches(0.5),
                 "→", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 11 — Daily 자동화
# ============================================================
s = add_blank_slide()
add_header_bar(s, 11, TOTAL, "AUTOMATION", "Daily 자동화 파이프라인")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6),
         "GitHub Actions cron으로 매일 KST 09:00경 자동 발송. 5중 cron + dedup 가드.",
         font_size=16, color=GRAY)

flow = [
    ("Cron 5개", "KST 08:47\n09:02 / 09:17\n09:32 / 09:52", ACCENT),
    ("Dedup 가드", "오늘 KST 성공\n이미 있으면\nskip", ACCENT2),
    ("Report 생성", "meta_jp_daily.py\nKPI + Off Rule\nBEST/WORST + 발견", GREEN),
    ("Gemini Audit", "본문 숫자 정합성\n80점+ PASS\n시에만 발송", GRAY),
    ("Gmail 발송", "se.heo@orbiters.co.kr\n매일 09:00경\n도착", NAVY),
]
y0 = Inches(3.0)
box_w = Inches(2.3)
gap = Inches(0.15)
total_w = box_w * 5 + gap * 4
start_x = (SLIDE_W - total_w) // 2

for i, (title, body, color) in enumerate(flow):
    x = start_x + (box_w + gap) * i
    add_rect(s, x, y0, box_w, Inches(3.2), fill=LIGHT)
    add_rect(s, x, y0, box_w, Inches(0.6), fill=color)
    add_text(s, x, y0 + Inches(0.1), box_w, Inches(0.5),
             title, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.0), box_w, Inches(2.0),
             body, font_size=12, color=NAVY, align=PP_ALIGN.CENTER)
    # 화살표
    if i < 4:
        add_text(s, x + box_w, y0 + Inches(1.3), gap, Inches(0.5),
                 "▶", font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 12 — Industry Consensus
# ============================================================
s = add_blank_slide()
add_header_bar(s, 12, TOTAL, "BENCHMARK", "Industry Consensus")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6),
         "광고주들이 본인 풀에서 어디를 winner / kill로 자르는지 — 3개 출처 일치.",
         font_size=16, color=GRAY)

# 출처 3개 박스
sources = [
    ("Motion 2026", "550K ads · $1.3B", "5% winner", "50% kill"),
    ("Alex Neiman", "DTC 6 brands · 1,847", "5.1% winner", "Bottom 60% kill"),
    ("coinis", "DTC framework", "Day 14+ scale", "3× CPA + 0 conv kill"),
]
y0 = Inches(2.8)
for i, (name, scale, winner, kill) in enumerate(sources):
    x = Inches(0.7 + i * 4.3)
    add_rect(s, x, y0, Inches(4.0), Inches(2.5), fill=LIGHT)
    add_text(s, x, y0 + Inches(0.2), Inches(4.0), Inches(0.5),
             name, font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(0.8), Inches(4.0), Inches(0.4),
             scale, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.4), Inches(4.0), Inches(0.4),
             winner, font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.9), Inches(4.0), Inches(0.4),
             kill, font_size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 하단: 우리 적용
add_rect(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.2), fill=NAVY)
add_text(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.5),
         "우리 17개 풀 적용 (5/15 합의)",
         font_size=14, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(6.4), Inches(11), Inches(0.5),
         "Top 10% (S 후보) + Bottom 20% (Off 후보, spend ≥ ₩30K 필터) — 풀 size 보정해 보수 설정",
         font_size=13, color=WHITE)


# ============================================================
# Slide 13 — 향후 과제 + 마무리
# ============================================================
s = add_blank_slide()
add_header_bar(s, 13, TOTAL, "ROADMAP", "향후 과제")

todos = [
    ("이번 달", "Meta × Rakuten 매출 매칭", "CTR만으로 효율 미검증 — ROI 가시화", ACCENT),
    ("다음 달", "S 후보 인플루언서 자동 발굴", "P90+ Top performer 풀 확대", ACCENT2),
    ("다음 달", "Frequency ≥ 2.3 자동 alert", "Ad fatigue 선제 대응", GREEN),
    ("6주 후", "PAID/GIFT KPI 차등 평가", "회수 압력 차이 반영", GRAY),
]
y0 = Inches(2.0)
for i, (when, title, body, color) in enumerate(todos):
    y = y0 + Inches(i * 1.0)
    add_rect(s, Inches(0.7), y, Inches(11.9), Inches(0.9), fill=LIGHT)
    add_rect(s, Inches(0.7), y, Inches(0.15), Inches(0.9), fill=color)
    add_text(s, Inches(1.1), y + Inches(0.25), Inches(1.8), Inches(0.5),
             when, font_size=14, bold=True, color=color)
    add_text(s, Inches(3.0), y + Inches(0.15), Inches(4.5), Inches(0.5),
             title, font_size=17, bold=True, color=NAVY)
    add_text(s, Inches(3.0), y + Inches(0.5), Inches(8.5), Inches(0.4),
             body, font_size=12, color=GRAY)

# 마무리 메시지
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         "본 기준은 살아있는 문서. 6주 운영 후 풀 분포 percentile로 임계 확정.",
         font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
         "변경 사유는 git 커밋 + 메모리에 기록.",
         font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


prs.save(str(OUT))
print(f"[OK] saved {OUT} ({OUT.stat().st_size:,} bytes)")
