"""
PPT 기획안: Grosmimi SG (DXOCQU0ERo2) 포맷 그대로 따라하기 — JP 버전
Reference: https://www.instagram.com/p/DXOCQU0ERo2/
원본 3슬라이드:
  1P: "What's So Special About Our Feeding Bottle? Let's find out!"
  2P: "So Thoughtfully Designed!" — 8 features 나열
  3P: "Available in 2 sizes!" — 200ml vs 300ml 비교 (Cap/Ring/Teat/Body)

JP 버전 제품 치환:
  PPSU ストローマグ 200ml (Stage 1) + ワンタッチ 300ml (Stage 2)
1080x1080 px (1:1)
"""
import io
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

SIDE = Emu(10287000)

PINK = RGBColor(0xFF, 0x6B, 0x9D)
LIGHT_PINK = RGBColor(0xFF, 0xD6, 0xE0)
IVORY = RGBColor(0xFF, 0xF5, 0xF7)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xFF, 0x8A, 0x7A)
GOLD = RGBColor(0xC9, 0xA3, 0x3A)
MINT = RGBColor(0xA8, 0xD8, 0xC9)

FONT_JP = "Yu Gothic UI"


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SIDE, SIDE)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *, size=40, bold=False,
             color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_JP):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_label(slide, left, top, width, height, text, *, bg=PINK, fg=WHITE, size=24):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(100000)
    tf.margin_top = tf.margin_bottom = Emu(50000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape


def add_placeholder_photo(slide, left, top, width, height, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_PINK
    shape.line.color.rgb = PINK
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[{label}]"
    run.font.name = FONT_JP
    run.font.size = Pt(20)
    run.font.color.rgb = PINK
    run.font.bold = True
    return shape


def add_brand_header(slide):
    """상단 GROSMIMI 워드마크 — 원본 3슬라이드 공통"""
    add_text(slide, 0, 500000, SIDE, 500000, "GROSMIMI",
             size=28, bold=True, color=PINK)


def slide1_cover(prs):
    """1P: What's So Special About Our Mug? Let's find out!"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)
    add_brand_header(s)

    # 메인 카피 (원본: "What's So Special About Our Feeding Bottle?")
    add_text(s, 300000, 1600000, 9687000, 900000, "グロミミのマグ、",
             size=44, bold=True, color=DARK)
    add_text(s, 300000, 2400000, 9687000, 1200000, "何がそんなに特別なの？",
             size=56, bold=True, color=PINK)

    # 서브 카피 (원본: "Let's find out!")
    add_text(s, 300000, 3700000, 9687000, 700000, "その理由、見てみましょう 👀",
             size=32, bold=False, color=CORAL)

    # 두 제품 placeholder
    add_placeholder_photo(s, 700000, 4800000, 4000000, 4200000,
                          "PPSU ストローマグ\n200ml\n실사진")
    add_placeholder_photo(s, 5500000, 4800000, 4000000, 4200000,
                          "ワンタッチ\n300ml\n실사진")

    add_text(s, 0, 9600000, SIDE, 400000, "1/4", size=16, color=GRAY)


def slide2_features(prs):
    """2P: So Thoughtfully Designed! — 8 features"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)
    add_brand_header(s)

    # 메인 타이틀 (원본: "So Thoughtfully Designed!")
    add_text(s, 300000, 1300000, 9687000, 700000, "すみずみまで",
             size=40, bold=True, color=DARK)
    add_text(s, 300000, 2000000, 9687000, 1000000, "こだわりの設計 ✨",
             size=56, bold=True, color=PINK)

    # 8 features (원본 구조 그대로 JP화)
    features = [
        ("医療グレードPPSU素材", "耐熱・耐久性抜群"),
        ("BPAフリー／無毒", "赤ちゃんも安心"),
        ("漏れにくい密閉設計", "カバン入れてもOK"),
        ("独自トライアングルテスト合格", "信頼できる品質"),
        ("人間工学に基づくスリムボディ", "小さな手にもフィット"),
        ("見やすい目盛り", "量がひと目でわかる"),
        ("軽量・割れにくい素材", "落としても安心"),
        ("MADE IN KOREA", "韓国No.1 ベビーマグ"),
    ]

    col_w = 4600000
    row_h = 650000
    left_col = 400000
    right_col = 5300000
    top_start = 3400000

    for i, (title, desc) in enumerate(features):
        row = i // 2
        col_x = left_col if i % 2 == 0 else right_col
        y = top_start + row * (row_h + 100000)

        # 체크 + 타이틀
        add_text(s, col_x, y, col_w, 400000,
                 f"✔ {title}",
                 size=20, bold=True, color=DARK,
                 align=PP_ALIGN.LEFT)
        # 서브 설명
        add_text(s, col_x + 250000, y + 380000, col_w, 300000,
                 desc,
                 size=14, bold=False, color=GRAY,
                 align=PP_ALIGN.LEFT)

    add_text(s, 0, 9600000, SIDE, 400000, "2/4", size=16, color=GRAY)


def slide3_sizes(prs):
    """3P: Available in 2 sizes! — PPSU 200ml vs ワンタッチ 300ml"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)
    add_brand_header(s)

    # 메인 타이틀 (원본: "Available in 2 sizes!")
    add_text(s, 300000, 1300000, 9687000, 700000, "成長に合わせて、",
             size=36, bold=True, color=DARK)
    add_text(s, 300000, 2000000, 9687000, 1000000, "2タイプから選べる！",
             size=52, bold=True, color=PINK)

    # 좌측 컬럼: PPSU 200ml
    col_left_x = 500000
    col_right_x = 5287000
    col_w = 4500000
    top_y = 3300000

    # 라벨
    add_label(s, col_left_x, top_y, col_w, 500000,
              "Stage 1 ｜ 6ヶ月〜", bg=PINK, fg=WHITE, size=20)

    # 제품명
    add_text(s, col_left_x, top_y + 650000, col_w, 500000,
             "PPSU ストローマグ 200ml",
             size=22, bold=True, color=DARK)

    # 제품 placeholder
    add_placeholder_photo(s, col_left_x + 500000, top_y + 1200000,
                          3500000, 3500000, "PPSU 200ml\n실사진")

    # 스펙 (원본 구조: Cap / Ring / Teat / Body)
    spec_y = top_y + 4900000
    specs_left = [
        ("キャップ", "ストローキャップ"),
        ("リング", "ソフトリング"),
        ("ストロー", "Stage 1 (柔らか)"),
        ("ボディ", "PPSU"),
    ]
    for i, (k, v) in enumerate(specs_left):
        add_text(s, col_left_x, spec_y + i * 300000, col_w, 280000,
                 f"・{k}：{v}",
                 size=15, bold=False, color=DARK, align=PP_ALIGN.LEFT)

    # 우측 컬럼: ワンタッチ 300ml
    add_label(s, col_right_x, top_y, col_w, 500000,
              "Stage 2 ｜ 12ヶ月〜", bg=CORAL, fg=WHITE, size=20)

    add_text(s, col_right_x, top_y + 650000, col_w, 500000,
             "ワンタッチ 300ml",
             size=22, bold=True, color=DARK)

    add_placeholder_photo(s, col_right_x + 500000, top_y + 1200000,
                          3500000, 3500000, "ワンタッチ 300ml\n실사진")

    specs_right = [
        ("キャップ", "ワンタッチ開閉"),
        ("リング", "耐久リング"),
        ("ストロー", "Stage 2 (しっかり)"),
        ("ボディ", "PPSU"),
    ]
    for i, (k, v) in enumerate(specs_right):
        add_text(s, col_right_x, spec_y + i * 300000, col_w, 280000,
                 f"・{k}：{v}",
                 size=15, bold=False, color=DARK, align=PP_ALIGN.LEFT)

    # 하단 카피 (원본: "Suitable for newborns to babies with larger appetite")
    add_text(s, 300000, 9150000, 9687000, 400000,
             "はじめての一本から、活動的な時期まで。",
             size=18, bold=True, color=DARK)

    add_text(s, 0, 9600000, SIDE, 400000, "3/4 · @grosmimi_japan",
             size=16, color=GRAY)


def slide4_cta(prs):
    """4P: CTA — 어그로 캐치프레이즈 + 프로필 팔로우 유도 (원본에 없는 추가 슬라이드)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)
    add_brand_header(s)

    # 어그로 후크 (상단)
    add_text(s, 300000, 1400000, 9687000, 900000,
             "「結局、どれがいいの…？」",
             size=52, bold=True, color=DARK)

    # 답 유도
    add_text(s, 300000, 2500000, 9687000, 700000,
             "その答え、プロフィールに全部あります📌",
             size=32, bold=True, color=PINK)

    # 중앙 placeholder (제품 라인업 이미지 or 인스타 아이콘)
    add_placeholder_photo(s, 2143000, 3700000, 6000000, 3500000,
                          "@grosmimi_japan\nプロフィール画面\nor 3제품 라인업 실사진")

    # 팔로우 버튼 스타일 CTA 박스
    cta = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(1500000), Emu(7500000), Emu(7287000), Emu(900000))
    cta.fill.solid()
    cta.fill.fore_color.rgb = PINK
    cta.line.fill.background()
    tf = cta.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "👆 @grosmimi_japan をフォロー"
    run.font.name = FONT_JP
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # 보조 카피
    add_text(s, 300000, 8700000, 9687000, 500000,
             "最新情報 & 限定クーポン、逃さず受け取ろう🎁",
             size=22, bold=False, color=DARK)

    add_text(s, 0, 9600000, SIDE, 400000, "4/4 · @grosmimi_japan",
             size=16, color=GRAY)


def slide5_caption(prs):
    """캡션 참조 슬라이드"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)

    add_text(s, 500000, 400000, 9287000, 600000, "【IG キャプション】",
             size=28, bold=True, color=PINK, align=PP_ALIGN.LEFT)

    caption = (
        "こんなにマグの種類がある中で、なぜグロミミ？🤔\n\n"
        "韓国生まれの受賞歴あるマグ🇰🇷\n"
        "医療グレードPPSU素材＆独自のトライアングルテストもクリア✨\n\n"
        "そして何より、お子さまの成長にあわせて使い分けできること🥰\n"
        "👉 はじめては『PPSUストローマグ 200ml』（6ヶ月〜）\n"
        "👉 自分で飲めるようになったら『ワンタッチ 300ml』へ\n"
        "同じシリーズだから、買い替えもラク🍃\n\n"
        "毎日のマグ、はじめてのマグはグロミミで🤍\n"
        "詳しくは @grosmimi_japan のプロフィールから\n\n"
        "#グロミミ #grosmimi #ストローマグ #スマートマグ "
        "#赤ちゃんのいる暮らし #育児グッズ #出産準備"
    )

    tb = s.shapes.add_textbox(Emu(500000), Emu(1200000), Emu(9287000), Emu(8000000))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(caption.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = FONT_JP
        run.font.size = Pt(20)
        run.font.color.rgb = DARK


def main():
    prs = Presentation()
    prs.slide_width = SIDE
    prs.slide_height = SIDE

    slide1_cover(prs)
    slide2_features(prs)
    slide3_sizes(prs)
    slide4_cta(prs)
    slide5_caption(prs)

    out_dir = Path(r"C:\Users\orbit\Desktop\s\요청하신 자료\인스타그램 기획안")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "なぜグロミミ_SG포맷_PPSU200_ワンタッチ300.pptx"
    prs.save(out)
    print(f"[OK] PPT saved: {out}")


if __name__ == "__main__":
    main()
