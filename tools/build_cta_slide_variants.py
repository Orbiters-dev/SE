"""
CTA 슬라이드 대안 5종 (미니멀 버전)
타입: 호기심/혜택/저장/공감/DM — 각 1장씩

디자인 원칙:
- 텍스트 최소 (메인 한 줄 + 보조 한 줄)
- 여백 많이
- 핑크 포인트 컬러로 포커스
- 1080x1080 (1:1)
"""
import io
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

SIDE = Emu(10287000)

PINK = RGBColor(0xFF, 0x6B, 0x9D)
LIGHT_PINK = RGBColor(0xFF, 0xE4, 0xEC)
IVORY = RGBColor(0xFF, 0xF8, 0xFA)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xFF, 0x8A, 0x7A)
MINT = RGBColor(0xB5, 0xE0, 0xD0)

FONT_JP = "Yu Gothic UI"


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SIDE, SIDE)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_text(slide, left, top, width, height, text, *, size=40, bold=False,
             color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_pill(slide, left, top, width, height, text, *, bg=PINK, fg=WHITE, size=20):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg


def add_dm_footer(slide):
    """하단 댓글 훅 + 쿠폰 DM 푸터 (ManyChat comment→DM 자동 응답용)"""
    # 배경 띠
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, Emu(8400000), SIDE, Emu(1000000))
    band.fill.solid()
    band.fill.fore_color.rgb = DARK
    band.line.fill.background()

    # 메인 훅 (댓글 유도)
    add_text(slide, 400000, 8500000, 9487000, 450000,
             "💬「マグ」とコメントしてね",
             size=22, bold=True, color=WHITE)

    # 서브 (혜택 — 댓글→DM 자동 전송)
    add_text(slide, 400000, 8950000, 9487000, 400000,
             "楽天の限定クーポン、DMで自動お届け🎁",
             size=16, bold=False, color=LIGHT_PINK)

    # 하단 핸들
    add_text(slide, 0, 9500000, SIDE, 350000, "@grosmimi_japan",
             size=14, bold=True, color=GRAY)


def add_handle(slide):
    """호환용 alias — 새 푸터 사용"""
    add_dm_footer(slide)


# ── 대안 1: 실수 고백 ────────────────────────────────
def variant_confession(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    # 태그
    add_pill(s, 3943000, 1600000, 2400000, 500000,
             "📩 3ヶ月前の私へ", bg=PINK, fg=WHITE, size=16)

    # 메인 (따옴표 안 후크)
    add_text(s, 400000, 3200000, 9487000, 900000,
             "「ストローマグ、",
             size=42, bold=True, color=DARK)
    add_text(s, 400000, 4100000, 9487000, 1200000,
             "どれも同じでしょ」",
             size=54, bold=True, color=PINK)

    # 서브 (펀치라인)
    add_text(s, 400000, 5700000, 9487000, 600000,
             "…って思ってた3ヶ月前の私、出ておいで🫣",
             size=22, bold=False, color=DARK)

    add_handle(s)


# ── 대안 2: 가격 쇼크 ────────────────────────────────
def variant_price_shock(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT_PINK)

    # 상단 라벨
    add_pill(s, 3943000, 1600000, 2400000, 500000,
             "👀 ぶっちゃけ", bg=WHITE, fg=PINK, size=16)

    # 가격 vs 가격
    add_text(s, 400000, 3000000, 9487000, 1100000,
             "1,000円 vs 3,000円。",
             size=56, bold=True, color=DARK)

    # 메인 질문
    add_text(s, 400000, 4400000, 9487000, 1200000,
             "違い、何？👀",
             size=72, bold=True, color=PINK)

    # CTA
    add_text(s, 400000, 6400000, 9487000, 600000,
             "答え、全部プロフィールに書いた📌",
             size=22, bold=False, color=DARK)

    add_handle(s)


# ── 대안 3: 숨긴 진실 (9할 어그로) ───────────────────
def variant_hidden_truth(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    # 아이콘/라벨
    add_pill(s, 3943000, 1600000, 2400000, 500000,
             "🤫 ここだけの話", bg=DARK, fg=WHITE, size=16)

    # 큰 숫자 포커스
    add_text(s, 400000, 2900000, 9487000, 1400000,
             "9割のママが",
             size=48, bold=True, color=DARK)
    add_text(s, 400000, 4300000, 9487000, 1500000,
             "見落としてる。",
             size=62, bold=True, color=PINK)

    # 서브 (힌트)
    add_text(s, 400000, 6200000, 9487000, 700000,
             "マグ選びの落とし穴、話していい？",
             size=24, bold=False, color=DARK)

    add_handle(s)


# ── 대안 4: 남편 설득 (위트) ────────────────────────
def variant_husband(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    # 위치
    add_pill(s, 3543000, 1600000, 3200000, 500000,
             "😏 パパ説得用", bg=PINK, fg=WHITE, size=16)

    # 본문
    add_text(s, 400000, 3000000, 9487000, 900000,
             "「なんでこれ選んだの？」",
             size=38, bold=True, color=DARK)

    # 답
    add_text(s, 400000, 4200000, 9487000, 1400000,
             "って聞かれたら、",
             size=34, bold=False, color=DARK)

    add_text(s, 400000, 5300000, 9487000, 1200000,
             "コレ見せて😏",
             size=62, bold=True, color=PINK)

    add_handle(s)


# ── 대안 5: 사회 증명 (자랑) ────────────────────────
def variant_social_proof(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, IVORY)

    # 상단
    add_pill(s, 3943000, 1600000, 2400000, 500000,
             "✋ 体感値です", bg=CORAL, fg=WHITE, size=16)

    # 질문 인용
    add_text(s, 400000, 2900000, 9487000, 800000,
             "ママ友「それ、どこの？」",
             size=32, bold=True, color=DARK)

    # 확률 박스
    add_text(s, 400000, 4100000, 9487000, 1700000,
             "聞かれる確率",
             size=30, bold=False, color=DARK)
    add_text(s, 400000, 4900000, 9487000, 1600000,
             "90%",
             size=120, bold=True, color=PINK)

    # 서브
    add_text(s, 400000, 6800000, 9487000, 600000,
             "(※グロミミ持ってるママ調べ)",
             size=18, bold=False, color=GRAY)

    add_handle(s)


def main():
    prs = Presentation()
    prs.slide_width = SIDE
    prs.slide_height = SIDE

    variant_confession(prs)    # 1 실수 고백
    variant_price_shock(prs)   # 2 가격 쇼크
    variant_hidden_truth(prs)  # 3 숨긴 진실 (9할)
    variant_husband(prs)       # 4 남편 설득
    variant_social_proof(prs)  # 5 사회 증명 (90%)

    out_dir = Path(r"C:\Users\orbit\Desktop\s\요청하신 자료\인스타그램 기획안")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "CTA_대안_5종.pptx"
    prs.save(out)
    print(f"[OK] PPT saved: {out}")


if __name__ == "__main__":
    main()
